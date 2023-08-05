# -*- encoding: utf-8 -*-
"""
This package implements text replacement in Powerpoint files in pptx format.

The text is searched and replaced in all possible places while preserving the
original character's formatting.

Text replacement can be configured to leave certain slides untouched (by specifying
which slides should be processed), or to not touching text in tables, charts or
text frames in any of the shapes.

This package can be imported and the class python_pptx_text_replacer used directly
or it can be called as main and given parameters to define what needs to be done.
"""
from __future__ import print_function, unicode_literals

import os
import sys
import argparse
import re

if sys.version_info[0]==3:
    import collections
    import collections.abc

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Inches

__version__ = "v0.0.1"

class TextReplacer:
    """
    This class implements text replacement in Powerpoint files in pptx format.
    The text is searched and replaced in all possible places.
    """
 
    def __init__(self, presentation_file_name,
                 tables=True,
                 charts=True,
                 textframes=True,
                 slides=''):
        
        self._replacements = []
        self._collected_replacements = []
        self._presentation_file_name = self._ensure_unicode(presentation_file_name)
        if not os.path.exists(self._presentation_file_name):
            raise ValueError("Presentation file %s does not exist." % ( self._presentation_file_name ))
        self._presentation = Presentation(presentation_file_name)
        self._tables = tables
        self._charts = charts
        self._textframes = textframes
        slide_cnt = len(self._presentation.slides)
        if len(slides.strip())==0:
            self._slides = [ True ] * slide_cnt
        else:
            self._slides = [ False ] * slide_cnt
            for rr in re.split('\\s*,\\s*',slides.strip()):
                r = re.split('\\s*-\\s*',rr,maxsplit=3)
                low = None
                high = None
                if len(r)<=2:
                    try:
                        low = int(r[0])
                    except:
                        low = None
                    high = low
                    if len(r)==2:
                        if len(r[1])==0:
                            high = slide_cnt
                        else:
                            try:
                                high = int(r[1])
                            except:
                                high = None
                if low is None or high is None:
                    raise ValueError('Slide list (--slides "%s") is not a comma separated list of slide numbers (i.e. 1) or slide number ranges (i.e. 4-12)' % (slides))
                if low<1 or low>slide_cnt:
                    raise ValueError('Slide number %s in list (--slides "%s") is lower than 1 or bigger than the last slide number %s' % ( low, slides, slide_cnt ))
                if high<1 or high>slide_cnt:
                    raise ValueError('Slide number %s in list (--slides "%s") is lower than 1 or bigger than the last slide number %s' % ( high, slides, slide_cnt ))
                if low > high:
                    raise ValueError('Slide range %s in list (--slides "%s") is invalid.' % ( r, slides ))
                for i in range(low-1,high):
                    self._slides[i] = True


    def replace_text(self, replacements):
        self._replacements = list( (self._ensure_unicode(match),self._ensure_unicode(repl)) for (match,repl) in replacements )
        self._collected_replacements.extend(replacements)

        for repl in self._replacements:
            if len(repl)==0:
                raise ValueError("A match string can not be empty.")

        for i in range(0,len(self._replacements)-1):
            srch_i = self._replacements[i][0]
            repl_i = self._replacements[i][1]
            for j in range(i+1, len(self._replacements)):
                srch_j = self._replacements[j][0]
                repl_j = self._replacements[j][1]
                if repl_i.find(srch_j)>=0:
                    print("WARNING: Replacement string %s at index %s matches search string %s at index %s. This may produce unintended results due to chained replacements!" % ( repl_i, i, srch_j, j))
                if srch_j.find(srch_i)>=0:
                    if srch_j.replace(srch_i,repl_i) == repl_j:
                        print("WARNING: Match/Replacement ('%s','%s') at index %s is obsolete due to match/replacement ('%s','%s') at index %s" % (srch_j,repl_j,j,srch_i,repl_i,i))
                    else:
                        print("WARNING: Match/Replacement ('%s','%s') at index %s will never match due to match/replacement ('%s','%s') at index %s" % (srch_j,repl_j,j,srch_i,repl_i,i))

        # loop through all slides
        slide_idx = 0
        print("Presentation[%s]" % (self._presentation_file_name))
        for slide in self._presentation.slides:
            print("  Slide[%s, id=%s] with title '%s'" % ( slide_idx+1, slide.slide_id, "<no title>" if slide.shapes.title is None else slide.shapes.title.text ))
            if self._slides[slide_idx]:
                self._process_shapes(2, slide)
            else:
                print("    ... skipped")
            slide_idx += 1

    def write_presentation_to_file(self, presentation_output_file_name):
        self._presentation.save(presentation_output_file_name)

    def get_replacements(self):
        return self._collected_replacements

    def get_presentation_file_name(self):
        return self._presentation_file_name

    def _ensure_unicode(self, text):
        if isinstance(text,(str,bytes) if sys.version_info.major==2 else bytes):
            return text.decode('UTF-8')
        return text

    def _replace_text_in_text_frame(self, level, text_frame):
        for (match, replacement) in self._replacements:
            pos_in_text_frame = self._ensure_unicode(text_frame.text).find(match)
            if pos_in_text_frame < 0:
                print("%sTrying to match '%s' -> no match" % ( "  "*level, match ))
            while pos_in_text_frame>=0:
                print("%sTrying to match '%s' -> matched at %s" % ( "  "*level, match, pos_in_text_frame ))
                to_match = match
                to_replace = replacement
                paragraph_idx = 0
                pos_in_paragraph = pos_in_text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph_len = len(self._ensure_unicode(paragraph.text))
                    if pos_in_paragraph >= paragraph_len:
                        pos_in_paragraph -= paragraph_len+1 # +1 for the new-line-character
                    else:
                        # this is the paragraph that contains the beginning of the match
                        (to_match, to_replace) = self._replace_runs_text(level+1, paragraph_idx, paragraph.runs, pos_in_paragraph, to_match, to_replace)
                        if len(to_match) == 0: # are we done with this match
                            break;
                        pos_in_paragraph = 0
                    paragraph_idx += 1
                pos_in_text_frame = self._ensure_unicode(text_frame.text).find(match,pos_in_text_frame+len(replacement))

    def _save_font_configuration(self, font):
        saved = {}
        saved['name'] = font.name
        saved['size'] = font.size
        saved['bold'] = font.bold
        saved['italic'] = font.italic
        saved['underline'] = font.underline
        saved['color.type'] = font.color.type
        if font.color.type == MSO_COLOR_TYPE.SCHEME:
            saved['color.brightness'] = font.color.brightness
            saved['color.theme_color'] = font.color.theme_color
        elif font.color.type == MSO_COLOR_TYPE.RGB:
            saved['color.rgb'] = None if font.color.rgb is None else str(font.color.rgb)
        # saved['fill'] = font.fill
        # saved['language_id'] = font.language_id
        return saved

    def _restore_font_configuration(self, saved, font):
        font.name = saved['name']
        font.size = saved['size']
        font.bold = saved['bold']
        font.italic = saved['italic']
        font.underline = saved['underline']
        if saved['color.type'] == MSO_COLOR_TYPE.SCHEME:
            font.color.brightness = saved['color.brightness']
            font.color.theme_color = saved['color.theme_color']
        elif saved['color.type'] == MSO_COLOR_TYPE.RGB:
            if saved['color.rgb'] is not None:
                font.color.rgb = RGBColor.from_string(saved['color.rgb'])
            else:
                font.color.rgb = None
        # font.fill = saved['fill']
        # font.language_id = saved['language_id']

    def _replace_runs_text(self, level, paragraph_idx, runs, pos, match, replacement):
        cnt = len(runs)
        i = 0
        while i<cnt:
            olen = len(self._ensure_unicode(runs[i].text))
            if pos>=olen:
                pos -= olen # the relative position of our match in the next run's text
                i += 1      # and off to the next run
            else:
                # we found the run, where the match starts!
                to_match = match
                match_len = len(to_match)
                to_replace = replacement
                repl_len = len(to_replace)

                while i<cnt:
                    run = runs[i]
                    otext = self._ensure_unicode(run.text)
                    olen = len(otext)
                    if pos+match_len < olen:
                        # our match ends before the end of the text of this run therefore
                        # we put the rest of our replacement string here and we are done!
                        saved_font = self._save_font_configuration(run.font)
                        run.text = otext[0:pos]+to_replace+otext[pos+match_len:]
                        self._restore_font_configuration(saved_font, run.font)
                        print("%sRun[%s,%s]: '%s' -> '%s'" % ( "  "*level, paragraph_idx, i, otext, run.text ))
                        return ('','')
                    if pos+match_len == olen:
                        # our match ends together with the text of this run therefore
                        # we put the rest of our replacement string here and we are done!
                        saved_font = self._save_font_configuration(run.font)
                        run.text = otext[0:pos]+to_replace
                        self._restore_font_configuration(saved_font, run.font)
                        print("%sRun[%s,%s]: '%s' -> '%s'" % ( "  "*level, paragraph_idx, i, otext, run.text ))
                        return ('','')
                    # we still haven't found all of our original match string
                    # so we process what we have here and go on to the next run
                    part_match_len = olen-pos
                    ntext = otext[0:pos]
                    if repl_len <= part_match_len:
                        # we now found at least as many characters for our match string
                        # as we have replacement characters for it. Thus we use up the
                        # the rest of our replacement string here and will replace the
                        # remainder of the match with an empty string (which happens
                        # to happen in this exact same spot for the next run ;-))
                        ntext += to_replace
                        repl_len = 0
                        to_replace = ''
                    else:
                        # we have got some more match characters but still more
                        # replacement characters than match characters found 
                        ntext += to_replace[0:part_match_len]
                        to_replace = to_replace[part_match_len:]
                        repl_len -= part_match_len
                    print("%sRun[%s,%s]: '%s' -> '%s'" % ( "  "*level, paragraph_idx, i, otext, ntext ))
                    saved_font = self._save_font_configuration(run.font)
                    run.text = ntext            # save the new text to the run
                    self._restore_font_configuration(saved_font, run.font)
                    to_match = to_match[part_match_len:] # this is what is left to match
                    match_len -= part_match_len # this is the length of the match that is left
                    pos = 0                     # in the next run, we start at pos 0 with our match
                    i += 1                      # and off to the next run
                return (to_match, to_replace)
            

    def _process_text_frame(self, level, text_frame):
        print("%sTextFrame: '%s'" % ( "  "*level, text_frame.text ))
        paragraph_idx = 0
        for paragraph in text_frame.paragraphs:
            print("%sParagraph[%s]: '%s'" % ( "  "*(level+1), paragraph_idx, paragraph.text ))
            run_idx = 0
            for run in paragraph.runs:
                print("%sRun[%s,%s]: '%s'" % ( "  "*(level+2), paragraph_idx, run_idx, run.text ))
                run_idx += 1
            paragraph_idx += 1
        self._replace_text_in_text_frame(level+1,text_frame)

    def _process_shapes(self, level, shape_list_parent):
        for shape in shape_list_parent.shapes:
            print("%sShape[%s, id=%s, type=%s]" % ( "  "*level, shape_list_parent.shapes.index(shape), shape.shape_id, shape.shape_type ))
            if shape.has_text_frame:
                if self._textframes:
                    self._process_text_frame(level+1,shape.text_frame)
                else:
                    print("%s... skipped" % ("  "*(level+1)))
            if shape.has_table:
                table = shape.table
                row_cnt = len(table.rows)
                col_cnt = len(table.columns)
                print("%sTable[%s,%s]" % ( "  "*(level+1), row_cnt, col_cnt ) )
                if self._tables:
                    for row in range(0, row_cnt):
                        for col in range(0, col_cnt):
                            cell = table.cell(row,col)
                            print("%sCell[%s,%s]: '%s'" % ( "  "*(level+2), row, col, cell.text ))
                            self._process_text_frame(level+3, cell.text_frame)
                else:
                    print("%s... skipped" % ( "  "*(level+2)))
            if shape.shape_type==MSO_SHAPE_TYPE.GROUP:
                self._process_shapes(level+1, shape)
            if shape.has_chart:
                chart = shape.chart
                print("%sChart of type %s" % ( "  "*(level+1), chart.chart_type ) )
                if self._charts:
                    categories_changed = False
                    new_categories = []
                    category_idx = 0
                    for category in chart.plots[0].categories:
                        print("%sCategory[%s] '%s'" % ( "  "*(level+2), category_idx, category ))
                        for (match,replace) in self._replacements:
                            changed_category = category.replace(match,replace)
                            if changed_category == category:
                                print("%sReplacing '%s' -> no match" % ( "  "*(level+3), match ))
                            else:
                                print("%sReplacing '%s' -> changed to '%s'" % ( "  "*(level+3), match, changed_category ))
                                category = changed_category
                                categories_changed = True
                        new_categories.append(category)
                        category_idx += 1

                    if categories_changed:
                        new_chart_data = CategoryChartData()
                        new_chart_data.categories = new_categories
                        for series in chart.series:
                            new_chart_data.add_series(series.name,series.values)
                        chart.replace_data(new_chart_data)
                else:
                    print("%s... skipped" % ( "  "*(level+2)))

def main():
    copyleft = "python-pptx-text-replacer %s (c) Frank Schäckermann 2022" % __version__
    p = argparse.ArgumentParser(description=__doc__,
                                formatter_class=argparse.RawDescriptionHelpFormatter,
                                epilog="""
The parameters --match and --replace can be specified multiple times.
They are paired up in the order of their appearance.

The slide list given with --slides must be a comma-separated list of
slide numbers from 1 to the number of slides contained in the presentation
or slide number ranges of the kind '4-16'. If the second number is omitted,
like in '4-' the range includes everything from the slide identified by the
first number up to the last slide in the file.

%s
%s
""" % ( "="*len(copyleft), copyleft ) )
    p.add_argument('--match',   '-m',
                   action='append',
                   required=True,
                   dest='matches',
                   metavar='<match>',
                   help='the string to look for and to be replaced')
    p.add_argument('--replace', '-r',
                   action='append',
                   required=True,
                   dest='replacements',
                   metavar='<replacement>',
                   help="the replacement for all the matches' occurrences")
    p.add_argument('--input',   '-i',
                   action='store',
                   required=True,
                   metavar='<input file>',
                   help="the file to replace the text in")
    p.add_argument('--output',  '-o',
                   action='store',
                   required=True,
                   metavar='<output file>',
                   help="the file to write the changed presentation to")
    p.add_argument('--slides', '-s',
                   metavar='<list of slide numbers to process>',
                   action='store',
                   required=False,
                   default='',
                   help="A comma-separated list of slide numbers (1-based) to restrict processing to, i.e. '2,4,6-10'")
    p.add_argument('--text-frames',  '-f',
                   action='store_const',
                   dest='textframes',
                   const=True,
                   required=False,
                   default=True,
                   help="process text frames in any shape as well (default)")
    p.add_argument('--no-text-frames','-F',
                   action='store_const',
                   dest='charts',
                   const=False,
                   required=False,
                   default=True,
                   help="do not process any text frames in shapes")
    p.add_argument('--tables',  '-t',
                   action='store_const',
                   dest='tables',
                   const=True,
                   required=False,
                   default=True,
                   help="process tables as well (default)")
    p.add_argument('--no-tables','-T',
                   action='store_const',
                   dest='tables',
                   const=False,
                   required=False,
                   default=True,
                   help="do not process tables and their cells")
    p.add_argument('--charts',  '-c',
                   action='store_const',
                   dest='charts',
                   const=True,
                   required=False,
                   default=True,
                   help="process chart categories as well (default)")
    p.add_argument('--no-charts','-C',
                   action='store_const',
                   dest='charts',
                   const=False,
                   required=False,
                   default=True,
                   help="do not process charts and their categories")

    ns = p.parse_args(sys.argv[1:])

    if len(ns.matches) != len(ns.replacements):
        print("There must be as many match-strings (-m) as there are replacement-strings (-r)", file=sys.stderr)
        return 1

    try:
        replacer = TextReplacer(ns.input,
                                tables=ns.tables,
                                charts=ns.charts,
                                textframes=ns.textframes,
                                slides=ns.slides)
        replacements = []
        for m in range(0,len(ns.matches)):
            replacements.append( ( ns.matches[m], ns.replacements[m] ) )
        
        replacer.replace_text(replacements)
        replacer.write_presentation_to_file(ns.output)

        return 0
    except ValueError as err:
        print(str(err.args[0]), file=sys.stderr)
        return 1

if __name__ == '__main__':
    sys.exit(main())
